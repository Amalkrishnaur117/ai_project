import os
from gtts import gTTS
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips

def create_audio_from_text(text, filename):
    tts = gTTS(text=text, lang='en')
    tts.save(filename)

def create_slide_image(text, slide_number):
    from PIL import Image, ImageDraw, ImageFont

    # Create an image with white background
    img = Image.new('RGB', (800, 600), color='white')
    d = ImageDraw.Draw(img)

    # You may need to adjust the font path
    font_path = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
    font = ImageFont.truetype(font_path, 40)

    # Draw text in the center
    text_width, text_height = d.textsize(text, font=font)
    d.text(((800 - text_width) / 2, (600 - text_height) / 2), text, fill="black", font=font)

    img.save(f"slide_{slide_number}.png")

def main(texts):
    audio_files = []
    slide_files = []

    for i, text in enumerate(texts):
        # Create slide images
        create_slide_image(text, i + 1)
        slide_files.append(f"slide_{i + 1}.png")

        # Create audio files
        audio_file = f"audio_{i + 1}.mp3"
        create_audio_from_text(text, audio_file)
        audio_files.append(audio_file)

    # Create video clips
    video_clips = []
    for slide_file, audio_file in zip(slide_files, audio_files):
        # Load slide and audio
        slide_clip = ImageClip(slide_file).set_duration(5)  # 5 seconds per slide
        audio_clip = AudioFileClip(audio_file)

        # Set audio to the slide
        slide_clip = slide_clip.set_audio(audio_clip)
        video_clips.append(slide_clip)

    # Concatenate all video clips
    final_video = concatenate_videoclips(video_clips)
    final_video.write_videofile("final_video.mp4", fps=24)

    # Clean up temporary files
    for file in audio_files + slide_files:
        os.remove(file)

if __name__ == "__main__":
    # Example text input
    texts = [
        "Welcome to the video!",
        "This is the second slide.",
        "And here's the third slide.",
        "Thanks for watching!"
    ]
    main(texts)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize a Presentation object
presentation = Presentation()

# Create a single slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background color (light blue)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue

# Add a title shape
title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
title_frame = title_shape.text_frame
title_frame.text = "Key Points"

# Customize the title style
title = title_frame.paragraphs[0]
title.font.size = Pt(36)  # Larger font size
title.font.bold = True  # Bold text
title.font.color.rgb = RGBColor(255, 69, 0)  # Red-Orange color
title.alignment = PP_ALIGN.CENTER  # Center the title

# Add a content shape for bullet points
content_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4.5))
content_frame = content_shape.text_frame

# List of points to add
points = [
    "Point 1: Introduction to the topic",
    "Point 2: Importance of the topic",
    "Point 3: Key findings",
    "Point 4: Future implications",
    "Point 5: Conclusion"
]

# Add the first point and set it as the bullet point
p = content_frame.add_paragraph()  # Add a new paragraph for the first point
p.text = points[0]
p.space_after = Pt(14)  # Space after the point
p.font.size = Pt(20)     # Font size for the point
p.font.color.rgb = RGBColor(0, 0, 0)  # Black color for the text
p.bullet = True  # Set bullet for the first point

# Add the remaining points
for point in points[1:]:
    p = content_frame.add_paragraph()  # Add a new paragraph for each subsequent point
    p.text = point
    p.space_after = Pt(14)  # Space after each point
    p.font.size = Pt(20)     # Font size for each point
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black color for the text
    p.bullet = True  # Set bullet for each point

# Save the presentation
presentation.save("attractive_single_slide_presentation.pptx")