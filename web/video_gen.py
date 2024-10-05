from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from gtts import gTTS
from moviepy.editor import VideoFileClip, AudioFileClip

# Step 1: Create the PowerPoint presentation
def create_presentation():
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)

    # Set the slide background color (light blue)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue

    # Add a title shape
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    title_frame = title_shape.text_frame
    title_frame.text = "Key Points"

    # Customize the title style
    title = title_frame.paragraphs[0]
    title.font.size = Pt(36)  # Larger font size
    title.font.bold = True  # Bold text
    title.font.color.rgb = RGBColor(255, 69, 0)  # Red-Orange color
    title.alignment = PP_ALIGN.CENTER  # Center the title

    # Add a content shape for bullet points
    content_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4.5))
    content_frame = content_shape.text_frame

    # List of points to add
    points = [
        "Point 1: Introduction to the topic",
        "Point 2: Importance of the topic",
        "Point 3: Key findings",
        "Point 4: Future implications",
        "Point 5: Conclusion"
    ]

    # Add the first point and set it as the bullet point
    p = content_frame.add_paragraph()  # Add a new paragraph for the first point
    p.text = points[0]
    p.space_after = Pt(14)  # Space after the point
    p.font.size = Pt(20)     # Font size for the point
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black color for the text
    p.bullet = True  # Set bullet for the first point

    # Add the remaining points
    for point in points[1:]:
        p = content_frame.add_paragraph()  # Add a new paragraph for each subsequent point
        p.text = point
        p.space_after = Pt(14)  # Space after each point
        p.font.size = Pt(20)     # Font size for each point
        p.font.color.rgb = RGBColor(0, 0, 0)  # Black color for the text
        p.bullet = True  # Set bullet for each point

    # Save the presentation
    presentation.save("attractive_single_slide_presentation.pptx")

# Step 2: Create audio from script
def create_audio(script):
    audio = gTTS(text=script, lang='en')
    audio_file = "presentation_audio.mp3"
    audio.save(audio_file)
    return audio_file

# Step 3: Combine audio with video
def combine_audio_video(video_file, audio_file):
    video_clip = VideoFileClip(video_file)
    audio_clip = AudioFileClip(audio_file)

    # Set the audio of the video
    final_clip = video_clip.set_audio(audio_clip)
    final_output = "final_presentation_video.mp4"
    final_clip.write_videofile(final_output, codec='libx264', audio_codec='aac')

    # Clean up
    video_clip.close()
    audio_clip.close()

# Main execution
if __name__ == "__main__":
    # Step 1: Create the PowerPoint presentation
    create_presentation()

    # Step 2: Define your script
    script = """
    Welcome to this presentation on Key Points.
    First, we will discuss the introduction to the topic.
    Next, we will highlight the importance of the topic.
    Then, we will cover key findings and future implications.
    Finally, we will conclude our presentation. Thank you for your attention!
    """

    audio_file = create_audio(script)

    # You will need to convert the PowerPoint to a video separately.
    # For demonstration purposes, we'll assume you have a video file named 'presentation_video.mp4'
    video_file = "presentation_video.mp4"  # Replace with your actual video file path

    # Step 3: Combine audio with video
    combine_audio_video(video_file, audio_file)

    print("Presentation video created successfully!")
