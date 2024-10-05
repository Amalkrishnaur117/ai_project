import os
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from gtts import gTTS
from moviepy.editor import ImageSequenceClip, AudioFileClip
from PIL import Image

def create_presentation(filename="presentation.pptx"):
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)

    # Background color (light blue)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    title_frame = title_shape.text_frame
    title_frame.text = "Key Points"
    title = title_frame.paragraphs[0]
    title.font.size = Pt(36)
    title.font.bold = True
    title.font.color.rgb = RGBColor(255, 69, 0)  # Red-Orange
    title.alignment = PP_ALIGN.CENTER

    # Add bullet points
    content_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4.5))
    content_frame = content_shape.text_frame
    points = [
        "Point 1: Introduction to the topic",
        "Point 2: Importance of the topic",
        "Point 3: Key findings",
        "Point 4: Future implications",
        "Point 5: Conclusion"
    ]
    for point in points:
        p = content_frame.add_paragraph()
        p.text = point
        p.space_after = Pt(14)
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.bullet = True

    presentation.save(filename)

def create_audio(script, audio_filename="presentation_audio.mp3"):
    audio = gTTS(text=script, lang='en')
    audio.save(audio_filename)

def convert_pptx_to_pdf(pptx_file, pdf_file="presentation.pdf"):
    subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", pptx_file, "--outdir", os.path.dirname(pdf_file)])

def convert_pdf_to_images(pdf_file, img_folder="slides"):
    if not os.path.exists(img_folder):
        os.makedirs(img_folder)

    with Image.open(pdf_file) as img:
        for i in range(img.n_frames):
            img.seek(i)
            img.save(os.path.join(img_folder, f"slide_{i + 1}.png"))

def create_video_from_images(img_folder="slides", output_file="presentation_video.mp4", fps=1):
    image_files = [os.path.join(img_folder, f) for f in sorted(os.listdir(img_folder)) if f.endswith('.png')]
    clip = ImageSequenceClip(image_files, fps=fps)
    clip.write_videofile(output_file, codec='libx264')

def combine_audio_video(video_file, audio_file, output_file="final_presentation_video.mp4"):
    video_clip = VideoFileClip(video_file)
    audio_clip = AudioFileClip(audio_file)
    final_clip = video_clip.set_audio(audio_clip)
    final_clip.write_videofile(output_file, codec='libx264', audio_codec='aac')

# Main execution
if __name__ == "__main__":
    # Create the PowerPoint presentation
    create_presentation("presentation.pptx")

    # Define script for audio
    script = """
    Welcome to this presentation on Key Points.
    First, we will discuss the introduction to the topic.
    Next, we will highlight the importance of the topic.
    Then, we will cover key findings and future implications.
    Finally, we will conclude our presentation. Thank you for your attention!
    """

    # Create audio file
    create_audio(script)

    # Convert PPTX to PDF
    convert_pptx_to_pdf("presentation.pptx")

    # Convert PDF pages to images
    convert_pdf_to_images("presentation.pdf")

    # Create video from images
    create_video_from_images()

    # Combine audio with video
    combine_audio_video("presentation_video.mp4", "presentation_audio.mp3")
