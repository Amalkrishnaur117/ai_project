from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(filename="attractive_single_slide_presentation.pptx"):
    # Initialize a Presentation object
    presentation = Presentation()

    # Create a single slide with a blank layout
    slide_layout = presentation.slide_layouts[5]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)

    # Set the slide background color (light blue)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue

    # Add a title shape
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    title_frame = title_shape.text_frame
    title_frame.text = "Key Points"

    # Customize the title style
    title = title_frame.paragraphs[0]
    title.font.size = Pt(36)  # Larger font size
    title.font.bold = True  # Bold text
    title.font.color.rgb = RGBColor(255, 69, 0)  # Red-Orange color
    title.alignment = PP_ALIGN.CENTER  # Center the title

    # Add a content shape for bullet points
    content_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4.5))
    content_frame = content_shape.text_frame

    # List of points to add
    points = [
        "Point 1: Introduction to the topic",
        "Point 2: Importance of the topic",
        "Point 3: Key findings",
        "Point 4: Future implications",
        "Point 5: Conclusion"
    ]

    # Add the first point and set it as the bullet point
    p = content_frame.add_paragraph()  # Add a new paragraph for the first point
    p.text = points[0]
    p.space_after = Pt(14)  # Space after the point
    p.font.size = Pt(20)     # Font size for the point
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black color for the text
    p.bullet = True  # Set bullet for the first point

    # Add the remaining points
    for point in points[1:]:
        p = content_frame.add_paragraph()  # Add a new paragraph for each subsequent point
        p.text = point
        p.space_after = Pt(14)  # Space after each point
        p.font.size = Pt(20)     # Font size for each point
        p.font.color.rgb = RGBColor(0, 0, 0)  # Black color for the text
        p.bullet = True  # Set bullet for each point

    # Save the presentation
    presentation.save(filename)

# Call the function to create the presentation
create_presentation()

